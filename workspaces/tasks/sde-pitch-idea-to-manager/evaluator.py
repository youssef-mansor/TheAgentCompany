import os
import logging
import pptx
from typing import List

from common import get_all_texts_from_slide, evaluate_with_llm, checkpoint
from scoring import Result, Checkpoint, bonus_for_completing_final

SLIDES_PATH = "/workspace/pitch.pptx"

@checkpoint
def grade_checkpoint1():
    return os.path.exists(SLIDES_PATH)

@checkpoint
def grade_checkpoint2(slides):
    """1st slide should contain the title <Idea: Integrate Sotopia into Openhands>"""
    if slides is None or len(slides.slides) < 1:
        return False
    try:
        intro_slide = slides.slides[0]
        text = get_all_texts_from_slide(intro_slide)
    except Exception as e:
        logging.error("Error reading slide 2: %s", e)
        return False

    return all(k in text for k in ["idea", "integrate", "sotopia", "into", "openhands"])

@checkpoint
def grade_checkpoint3(slides):
    """2nd slide should contain a short introduction to Sotopia"""
    if slides is None or len(slides.slides) < 2:
        return False
    try:
        intro_slide = slides.slides[1]
        text = get_all_texts_from_slide(intro_slide)
    except Exception as e:
        logging.error("Error reading slide 2: %s", e)
        return False

    return evaluate_with_llm(text, "the content contains a short introduction to Sotopia.\
        FYI, Sotopia is an open-ended social learning environment that allows agents to interact with each other and the environment. \
        The environment is designed to be a platform for evaluating and faciliating social intelligence in language agents. \
        The environment is designed to be open-ended, meaning that the environment can be easily extended to include new environments and new agents. \
        The environment is also designed to be scalable, meaning that the environment can be easily scaled to include a large number of agents and environments.")

@checkpoint
def grade_checkpoint4(slides):
    """3nd slide has the title <Social Tasks>"""
    if slides is None or len(slides.slides) < 3:
        return False
    try:
        idea_slide = slides.slides[2]
        text = get_all_texts_from_slide(idea_slide)
    except Exception as e:
        logging.error("Error reading slide 3: %s", e)
        return False

    return all(k in text for k in ["social", "task"])

@checkpoint
def grade_checkpoint5(slides):
    """4th slide should be Q&A"""
    if slides is None or len(slides.slides) < 4:
        return False
    try:
        idea_slide = slides.slides[3]
        text = get_all_texts_from_slide(idea_slide)
    except Exception as e:
        logging.error("Error reading slide 4: %s", e)
        return False

    return all(k in text for k in ["q&a"])

def grade_checkpoints(trajectory="") -> Result:
    checkpoint1 = Checkpoint(1, 1*int(grade_checkpoint1()))
    try:
        slides = pptx.Presentation(SLIDES_PATH)
    except Exception as e:
        logging.error("Error reading slides: %s", e)
        slides = None
    checkpoint2 = Checkpoint(1, 1*int(grade_checkpoint2(slides)))
    checkpoint3 = Checkpoint(1, 1*int(grade_checkpoint3(slides)))
    checkpoint4 = Checkpoint(1, 1*int(grade_checkpoint4(slides)))
    checkpoint5 = Checkpoint(1, 1*int(grade_checkpoint5(slides)))

    return Result([checkpoint1, checkpoint2, checkpoint3, checkpoint4, checkpoint5])
