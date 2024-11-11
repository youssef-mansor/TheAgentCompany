import os
import logging
from typing import List

import pptx

from common import create_rocketchat_client
from common import get_rocketchat_personal_chat_history
from common import get_all_texts_from_slide, grader
from scoring import Result, Checkpoint

# Initialize RocketChat client
rocket = create_rocketchat_client()

LEADERSHIP_NAMES = ["Sarah Johnson", "Mark Johnson", "Jessica Lee", "David Wong", "Chen Xinyi"]
LEADERSHIP_POSITIONS = ["CTO", "Sales Director", "Marketing Manager", "Finance Director", "Human Resources Manager"]
SLIDES_PATH = '/workspace/leadership_intro.pptx'


@grader
def grade_checkpoint1():
    return os.path.exists(SLIDES_PATH)


@grader
def grade_checkpoint2():
    return bool(get_rocketchat_personal_chat_history(rocket, "Chen Xinyi"))


@grader
def grade_checkpoint3():
    if not os.path.exists(SLIDES_PATH):
        return False

    try:
        slides = pptx.Presentation(SLIDES_PATH).slides
        return len(slides) == 5
    except Exception as e:
        logging.warning(f"Error opening slides: {e}")
        return False


@grader
def grade_checkpoint4():
    if not os.path.exists(SLIDES_PATH):
        return False

    try:
        slides = pptx.Presentation(SLIDES_PATH).slides
        for slide, name in zip(slides, LEADERSHIP_NAMES):
            slide_texts = get_all_texts_from_slide(slide)
            if name.lower() not in slide_texts:
                logging.warning(f"Name {name} not found in slide {slide_texts}")
                return False
    except Exception as e:
        logging.warning(f"Error getting slide texts: {e}")
        return False
    return True


@grader
def grade_checkpoint5():
    if not os.path.exists(SLIDES_PATH):
        return False

    try:
        slides = pptx.Presentation(SLIDES_PATH).slides
        for slide, position in zip(slides, LEADERSHIP_POSITIONS):
            slide_texts = get_all_texts_from_slide(slide)
            if position.lower() not in slide_texts:
                logging.warning(f"Position {position} not found in slide {slide_texts}")
                return False
    except Exception as e:
        logging.warning(f"Error getting slide texts: {e}")
        return False
    return True


def grade_checkpoints(trajectory=""):
    checkpoints: List[Checkpoint] = [
        Checkpoint(1, int(grade_checkpoint1())),
        Checkpoint(1, int(grade_checkpoint2())),
        Checkpoint(1, int(grade_checkpoint3())),
        Checkpoint(1, int(grade_checkpoint4())),
        Checkpoint(1, int(grade_checkpoint5())),
    ]
    result = Result(checkpoints)
    return result
