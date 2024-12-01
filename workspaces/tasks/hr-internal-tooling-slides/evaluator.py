import logging
from typing import List
from io import BytesIO

from pptx import Presentation

from common import create_rocketchat_client, check_file_in_owncloud_directory, get_binary_file_content_owncloud, evaluate_with_llm, evaluate_chat_history_with_llm, grader
from scoring import Result, Checkpoint
from config import OWNCLOUD_URL, ROCKETCHAT_URL, GITLAB_BASEURL, PLANE_BASEURL

# Initialize RocketChat client
rocket = create_rocketchat_client()

theme_keywords=["internal", "tooling", "training"]
rocketchat_channels = ["#general", "#kudos", "#engineering", "#help-desk", "#hr-announcements", "#product", 
                       "#project-ai", "#project-graphdb", "#project-lowcode", "#project-search", "#project-streamdb", 
                       "#project-webcrawler", "#random", "#sales-marketing", "#tech-talk"]
gitlab_repos = ["api-server", "bustub", "colly", "copilot-arena-server", "Documentation", "janusgraph", 
                "llama.cpp", "node-red", "openhands", "opensearch", "raft", "risingwave", "sotopia", "streamlit"]
folder_structure = ["Research", "Human Resource Team", "Financials", "Engineering", "Data Analysis", 
                    "Administrative Specialist", "Admin"]
plane_keywords = ["Total Open Tasks", "Pending Issues", "Most Issues Created", "Most Issues Closed"]

def get_title(slide):
    return slide.shapes.title.text if slide.shapes.title else ""

def get_content(slide):
    return slide.placeholders[1].text if len(slide.placeholders) > 1 else ""

def check_rocketchat(slide, required_channels, url):
    content = get_content(slide)
    channels_present = all(channel in content for channel in required_channels)
    check_default_channel = evaluate_with_llm(content, "default channels are #general and #kudos")
    return get_title(slide) == "RocketChat" and channels_present and url in content and check_default_channel

def check_gitlab(slide, url):
    content = get_content(slide)
    keywords = ["code collaboration", "project management", "issue tracking"]
    return get_title(slide) == "Gitlab" and all(keyword in content for keyword in keywords) and url in content

def check_Owncloud(slide, folder_structure, url):
    content = get_content(slide)
    
    folders_present = evaluate_with_llm(content, f"the content includes intros to all folders in below folders: {" ".join(folder_structure)}")
    return get_title(slide) == "Owncloud" and folders_present and url in content

def check_plane(slide, url):
    content = get_content(slide)
    analytics_intro = evaluate_with_llm(content, f"introduction to functionalities of plane analytics page, including following functions: {" ".join(plane_keywords)}")
    return get_title(slide) == "Plane" and analytics_intro and url in content

def check_content_margins(slides):
    for slide in slides:
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                return False
    return True

def check_markdown_formatting(slides):
    for slide in slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and ("**" in shape.text or "*" in shape.text):
                return False
    return True


@grader
def grade_checkpoint1():
    file_name = "Internal_Tooling_Training.pptx"
    dir_name = "Documents/Human%20Resources%20Team"
    if not check_file_in_owncloud_directory(file_name, dir_name):
        return 0

    content = get_binary_file_content_owncloud(file_name, dir_name)
    if not content:
        logging.warning("Failed to download content")
        return 0

    try:
        prs = Presentation(BytesIO(content))
    except Exception as e:
        logging.warning(f"Failed to read {file_name}: {e}")
        return 0
    
    slides = prs.slides

    TAC_content = get_binary_file_content_owncloud("TAC_overview.md", "Documents").decode("utf-8")

    if len(slides) != 7:
        return 0

    results = {
        "Title Check": any(keyword in get_title(slides[0]).lower() for keyword in theme_keywords),
        "TAC Overview": get_title(slides[1]) == "TAC Overview" and evaluate_with_llm(get_content(slides[1]), "the contents include introduction to " + TAC_content),
        "RocketChat": check_rocketchat(slides[2], rocketchat_channels, ROCKETCHAT_URL),
        "Gitlab Usage": check_gitlab(slides[3], GITLAB_BASEURL),
        "Gitlab Repos": get_title(slides[4]) == "Gitlab Repos" and all(repo in get_content(slides[4]) for repo in gitlab_repos),
        "Owncloud Folders": check_Owncloud(slides[5], folder_structure, OWNCLOUD_URL),
        "Plane Overview": check_plane(slides[6], PLANE_BASEURL),
    }

    for key, value in results.items():
        if not value:
            logging.warning(f"{key} check failed")

    return sum(results.values())


@grader
def grade_checkpoint2():
    file_name = "Internal_Tooling_Training.pptx"
    dir_name = "Documents/Human%20Resources%20Team"
    if not check_file_in_owncloud_directory(file_name, dir_name):
        return 0

    content = get_binary_file_content_owncloud(file_name, dir_name)
    if not content:
        logging.warning("Failed to download content")
        return 0

    try:
        prs = Presentation(BytesIO(content))
    except Exception as e:
        logging.warning(f"Failed to read {file_name}: {e}")
        return 0
    
    slides = prs.slides

    style_results = {
        "Content Margins": check_content_margins(slides),
        "No Markdown Markings": check_markdown_formatting(slides),
    }

    return sum(style_results.values())


@grader
def grade_checkpoint3():
    return evaluate_chat_history_with_llm(rocket, "Chen Xinyi", f"There is a message about review and inclusion of slides in the onboarding materials, and includes the link to the slides.")


def grade_checkpoints(trajectory=""):
    checkpoints: List[Checkpoint] = [
        Checkpoint(7, int(grade_checkpoint1())),
        Checkpoint(2, int(grade_checkpoint2())),
        Checkpoint(1, int(grade_checkpoint3())),
    ]
    result = Result(checkpoints)
    return result
