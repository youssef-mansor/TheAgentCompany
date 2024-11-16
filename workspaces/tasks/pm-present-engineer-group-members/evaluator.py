import os
import logging
import pptx
from typing import List

from common import get_all_texts_from_slide, evaluate_chat_history_with_llm, create_rocketchat_client, grader
from scoring import Result, Checkpoint, bonus_for_completing_final

SLIDES_PATH = "/workspace/engineers.pptx"
NPC_NAME = "chen_xinyi"

ENGINEERS = [
    ["Sarah", "Johnson", "CTO"],
    ["Li", "Ming", "Database Team Project Manager"],
    ["Zhang", "Wei", "Senior Software Engineer"],
    ["Wang", "Fang", "AI Researcher"],
    ["Mike", "Chen", "Senior Software Engineer"],
    ["Emily", "Zhou", "Software Engineer"]    
]

@grader
def grade_checkpoint1(slides):
    """Slides should contain information about all engineers"""
    if slides is None or len(slides.slides) < 1:
        return False
    text = ""
    for i in range(len(slides.slides)):
        try:
            text += get_all_texts_from_slide(slides.slides[i])
        except Exception as e:
            logging.error("Error reading slide 2: %s", e)
            return False
    logging.info("Full texts extracted from slides: %s", text)
    for engineer in ENGINEERS:
        if not all(k.lower() in text for k in engineer):
            logging.error("Engineer %s not found in slides", engineer)
            return False
    return True

@grader
def grade_checkpoint2():
    rocket_client = create_rocketchat_client()
    return evaluate_chat_history_with_llm(rocket_client, NPC_NAME, "introductory slides are finished")


def grade_checkpoints(trajectory="") -> Result:
    try:
        slides = pptx.Presentation(SLIDES_PATH)
    except Exception as e:
        logging.error("Error reading slides: %s", e)
        slides = None
    pass_checkpoint1 = grade_checkpoint1(slides)
    checkpoint1 = Checkpoint(2, 2*int(pass_checkpoint1))
    checkpoint2 = Checkpoint(1, 1*int(pass_checkpoint1 and grade_checkpoint2()))

    return Result([checkpoint1, checkpoint2])
